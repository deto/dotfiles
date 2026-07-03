from __future__ import annotations

import argparse
import io
import re
import sys
import time
from pathlib import Path

from google.auth.transport.requests import AuthorizedSession
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseUpload
from pptx import Presentation

from deck_drive import PPTX_MIME, download_pptx, drive_client


DEFAULT_DECK_SA = "deck-agent-sa@cloud-cli-481516.iam.gserviceaccount.com"
DEFAULT_TMP_FOLDER_ID = "1mVFIlE59dnMM0BGyvxfp-abUOA760_ko"  # .agent_tmp


def slugify(value: str) -> str:
    value = re.sub(r"[^A-Za-z0-9._-]+", "_", value.strip())
    value = re.sub(r"_+", "_", value).strip("._")
    return value or "drive_slide_preview"


def remove_all_but(prs: Presentation, keep_zero_index: int) -> None:
    slide_ids = prs.slides._sldIdLst
    for idx in reversed(range(len(prs.slides))):
        if idx != keep_zero_index:
            slide_ids.remove(slide_ids[idx])


def upload_preview_deck(
    drive,
    pptx_bytes: bytes,
    *,
    name: str,
    folder_id: str,
) -> dict:
    media = MediaIoBaseUpload(io.BytesIO(pptx_bytes), mimetype=PPTX_MIME, resumable=True)
    return (
        drive.files()
        .create(
            body={"name": name, "mimeType": PPTX_MIME, "parents": [folder_id]},
            media_body=media,
            supportsAllDrives=True,
            fields="id,name,webViewLink,thumbnailLink,hasThumbnail,thumbnailVersion",
        )
        .execute()
    )


def get_thumbnail_meta(drive, file_id: str) -> dict:
    return (
        drive.files()
        .get(
            fileId=file_id,
            supportsAllDrives=True,
            fields="id,name,webViewLink,thumbnailLink,hasThumbnail,thumbnailVersion",
        )
        .execute()
    )


def wait_for_thumbnail(
    drive,
    file_id: str,
    *,
    timeout_s: int,
    interval_s: int,
) -> dict:
    deadline = time.time() + timeout_s
    last_meta = {}
    while time.time() < deadline:
        last_meta = get_thumbnail_meta(drive, file_id)
        if last_meta.get("thumbnailLink"):
            return last_meta
        time.sleep(interval_s)
    raise RuntimeError(f"Drive did not generate thumbnailLink. Last metadata: {last_meta}")


def download_thumbnail(drive, thumbnail_link: str, out_path: Path, size: int) -> None:
    link = re.sub(r"=s\d+$", f"=s{size}", thumbnail_link)
    session = AuthorizedSession(drive._http.credentials)
    response = session.get(link, timeout=30)
    response.raise_for_status()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_bytes(response.content)


def trash_or_delete_file(drive, file_id: str) -> None:
    try:
        drive.files().update(
            fileId=file_id,
            body={"trashed": True},
            supportsAllDrives=True,
            fields="id,trashed",
        ).execute()
    except HttpError as update_error:
        try:
            drive.files().delete(fileId=file_id, supportsAllDrives=True).execute()
        except HttpError as delete_error:
            print(
                "warning: could not trash/delete temporary Drive file "
                f"{file_id}: update={update_error}; delete={delete_error}",
                file=sys.stderr,
            )


def deck_name(drive, file_id: str) -> str:
    meta = (
        drive.files()
        .get(fileId=file_id, supportsAllDrives=True, fields="name")
        .execute()
    )
    return meta.get("name", file_id)


def create_single_slide_pptx(
    drive,
    *,
    source_file_id: str,
    slide_number: int,
) -> bytes:
    source = download_pptx(drive, source_file_id)
    prs = Presentation(io.BytesIO(source))
    if not 1 <= slide_number <= len(prs.slides):
        raise ValueError(f"Slide {slide_number} not in deck with {len(prs.slides)} slides")

    remove_all_but(prs, slide_number - 1)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Render a Drive PPTX slide preview by uploading a temporary "
            "single-slide deck and downloading Drive's file thumbnail."
        )
    )
    parser.add_argument("--file-id", required=True, help="Source Drive PPTX file ID.")
    parser.add_argument(
        "--slide-number",
        required=True,
        type=int,
        help="1-based slide number to preview.",
    )
    parser.add_argument(
        "--out-dir",
        default="tmp",
        help="Local directory for the one-slide PPTX and thumbnail PNG.",
    )
    parser.add_argument(
        "--deck-sa",
        default=DEFAULT_DECK_SA,
        help="Service account to impersonate for Drive access.",
    )
    parser.add_argument(
        "--tmp-folder-id",
        default=DEFAULT_TMP_FOLDER_ID,
        help="Writable Drive folder for temporary preview decks. Defaults to .agent_tmp.",
    )
    parser.add_argument(
        "--name-prefix",
        default="drive_slide_preview",
        help="Prefix for local files and the temporary Drive deck name.",
    )
    parser.add_argument(
        "--thumbnail-size",
        default=1600,
        type=int,
        help="Requested Drive thumbnail size parameter.",
    )
    parser.add_argument(
        "--timeout-s",
        default=120,
        type=int,
        help="How long to wait for Drive thumbnail generation.",
    )
    parser.add_argument(
        "--interval-s",
        default=5,
        type=int,
        help="Polling interval while waiting for thumbnail generation.",
    )
    parser.add_argument(
        "--delete-drive-file",
        action="store_true",
        help="Trash the temporary Drive deck after downloading the thumbnail.",
    )
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = parse_args(sys.argv[1:] if argv is None else argv)
    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    drive = drive_client(args.deck_sa)
    source_name = deck_name(drive, args.file_id)
    base = slugify(f"{args.name_prefix}_{source_name}_slide_{args.slide_number}")
    local_pptx = out_dir / f"{base}.pptx"
    local_thumbnail = out_dir / f"{base}_thumbnail.png"
    drive_name = f"{base}.pptx"

    pptx_bytes = create_single_slide_pptx(
        drive,
        source_file_id=args.file_id,
        slide_number=args.slide_number,
    )
    local_pptx.write_bytes(pptx_bytes)

    created = upload_preview_deck(
        drive,
        pptx_bytes,
        name=drive_name,
        folder_id=args.tmp_folder_id,
    )
    try:
        meta = wait_for_thumbnail(
            drive,
            created["id"],
            timeout_s=args.timeout_s,
            interval_s=args.interval_s,
        )
        download_thumbnail(
            drive,
            meta["thumbnailLink"],
            local_thumbnail,
            args.thumbnail_size,
        )
    finally:
        if args.delete_drive_file:
            trash_or_delete_file(drive, created["id"])

    print(f"source_name: {source_name}")
    print(f"source_file_id: {args.file_id}")
    print(f"slide_number: {args.slide_number}")
    print(f"drive_file_id: {created['id']}")
    print(f"drive_name: {created['name']}")
    print(f"webViewLink: {meta.get('webViewLink')}")
    print(f"local_pptx: {local_pptx}")
    print(f"thumbnail: {local_thumbnail}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
