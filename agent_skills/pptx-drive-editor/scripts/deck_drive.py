"""
deck_drive.py — edit PowerPoint (.pptx) decks stored in a Google Shared Drive,
in place, using in-code service-account impersonation.

Why this shape:
  * No key file on disk and no clobbering the machine's default ADC — the
    base identity (your `gcloud auth application-default login`, or an attached
    SA on GCP) is impersonated up to the deck SA at runtime, in code.
  * Raw-OOXML editing (unzip -> edit slideN.xml -> repack) only rewrites the
    parts you touch, so template formatting, media, fonts, and relationships
    are preserved byte-for-byte and text stays editable.

Access model (two independent grants — keep them straight):
  * IAM:  base identity needs roles/iam.serviceAccountTokenCreator on the deck
          SA  ->  governs *who may act as* the SA.
  * Drive: deck SA must be a member of the Shared Drive / folder (Content
          Manager)  ->  governs *what content* the SA can read/overwrite.

Dependencies: google-api-python-client, google-auth
              (python-pptx only if you use insert_picture)
"""

from __future__ import annotations

import io
import re
import zipfile
from typing import Callable

import google.auth
from google.auth import impersonated_credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload, MediaIoBaseUpload

DRIVE_SCOPE = "https://www.googleapis.com/auth/drive"
PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument."
    "presentationml.presentation"
)
_SLIDE_RE = re.compile(r"ppt/slides/slide\d+\.xml")
_SHARED = {"supportsAllDrives": True}  # required for every Shared Drive call


# --------------------------------------------------------------------------- #
# Auth — impersonate the deck SA in code (no key file, no ADC clobber)
# --------------------------------------------------------------------------- #
def drive_client(deck_sa_email: str):
    """Build a Drive v3 client that acts as `deck_sa_email`.

    Resolves whatever base ADC exists, then mints a short-lived, Drive-scoped
    token for the deck SA. Nothing is written to disk.
    """
    source, _ = google.auth.default(
        scopes=["https://www.googleapis.com/auth/cloud-platform"]
    )
    creds = impersonated_credentials.Credentials(
        source_credentials=source,
        target_principal=deck_sa_email,
        target_scopes=[DRIVE_SCOPE],
    )
    return build("drive", "v3", credentials=creds, cache_discovery=False)


# --------------------------------------------------------------------------- #
# Drive I/O — Shared Drive aware
# --------------------------------------------------------------------------- #
def download_pptx(drive, file_id: str) -> bytes:
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, drive.files().get_media(fileId=file_id, **_SHARED))
    done = False
    while not done:
        _, done = dl.next_chunk()
    return buf.getvalue()


def upload_pptx_in_place(drive, file_id: str, data: bytes) -> dict:
    """Overwrite the SAME file. Keeps fileId -> links, sharing, comments survive."""
    media = MediaIoBaseUpload(io.BytesIO(data), mimetype=PPTX_MIME, resumable=True)
    return drive.files().update(fileId=file_id, media_body=media, **_SHARED).execute()


# --------------------------------------------------------------------------- #
# Raw-OOXML editing
# --------------------------------------------------------------------------- #
def read_slides(data: bytes) -> dict[str, str]:
    """Return {slide_part_name: xml_text} for every slide, in deck order.

    Use this to let the agent INSPECT slide XML before deciding what to change
    (which box is the title, where a placeholder lives, etc.).
    """
    zf = zipfile.ZipFile(io.BytesIO(data))
    names = sorted(
        (n for n in zf.namelist() if _SLIDE_RE.fullmatch(n)),
        key=lambda n: int(re.search(r"(\d+)", n.rsplit("/", 1)[1]).group(1)),
    )
    return {n: zf.read(n).decode("utf-8") for n in names}


def edit_slides(data: bytes, edit_fn: Callable[[str, str], str]) -> bytes:
    """Rewrite slide XML via edit_fn(part_name, xml) -> new_xml.

    Only slide parts are passed through edit_fn; every other part is copied
    byte-for-byte. Return the XML unchanged to leave a slide alone.
    """
    src = zipfile.ZipFile(io.BytesIO(data))
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            blob = src.read(item.filename)
            if _SLIDE_RE.fullmatch(item.filename):
                blob = edit_fn(item.filename, blob.decode("utf-8")).encode("utf-8")
            dst.writestr(item, blob)
    return out.getvalue()


def replace_tokens(data: bytes, mapping: dict[str, str]) -> bytes:
    """Convenience: literal token -> value across all slide XML.

    Reliable only when each token sits contiguously in ONE <a:t> run.
    PowerPoint sometimes splits a token across runs (autocorrect/spell-check);
    if a replacement silently no-ops, that's why — author template tokens in a
    single motion so they stay one run, or edit XML directly via edit_slides().
    """
    def esc(s: str) -> str:
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    escaped = {tok: esc(val) for tok, val in mapping.items()}

    def _fn(_name: str, xml: str) -> str:
        for tok, val in escaped.items():
            xml = xml.replace(tok, val)
        return xml

    return edit_slides(data, _fn)


# --------------------------------------------------------------------------- #
# Adding NEW pictures/charts — python-pptx is the pragmatic tool here.
# (Inserting media via raw OOXML means hand-managing content-types, rels, and
#  <p:pic> geometry; python-pptx does it correctly. Tradeoff: it rewrites more
#  of the package than a raw edit, so prefer raw-OOXML for text/structure.)
# --------------------------------------------------------------------------- #
def insert_picture(data: bytes, slide_index: int, image_path: str,
                   left_in: float = 1.0, top_in: float = 1.0,
                   width_in: float = 4.0) -> bytes:
    """Auto-place a figure (e.g. a matplotlib/plotly PNG) — no manual pasting."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(io.BytesIO(data))
    prs.slides[slide_index].shapes.add_picture(
        image_path, Inches(left_in), Inches(top_in), width=Inches(width_in)
    )
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# --------------------------------------------------------------------------- #
# End-to-end convenience
# --------------------------------------------------------------------------- #
def edit_deck_in_place(deck_sa_email: str, file_id: str,
                       edit_fn: Callable[[str, str], str]) -> dict:
    """download -> raw-OOXML edit -> overwrite same file. Returns Drive metadata."""
    drive = drive_client(deck_sa_email)
    data = download_pptx(drive, file_id)
    edited = edit_slides(data, edit_fn)
    return upload_pptx_in_place(drive, file_id, edited)
